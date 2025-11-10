# pdf_utils.py
from io import BytesIO
from typing import List
from pypdf import PdfReader  # pypdf>=4
from pptx import Presentation  # python-pptx

# We'll call OpenAI Vision from here lazily to avoid circular imports
def _vision_extract_text(image_bytes: bytes) -> str:
    try:
        from base64 import b64encode
        from llm import get_client, _supports_responses
        b64 = b64encode(image_bytes).decode("utf-8")
        data_url = f"data:image/png;base64,{b64}"  # works for jpg too

        client = get_client()
        # Use chat.completions vision input; widely supported
        messages = [{
            "role": "user",
            "content": [
                {"type": "text", "text": "Transcribe all legible text in reading order. If formulas/equations, retain them as LaTeX."},
                {"type": "image_url", "image_url": {"url": data_url}},
            ],
        }]
        try:
            # new client path
            resp = client.chat.completions.create(model="gpt-4o-mini", messages=messages, temperature=0.0)
            return resp.choices[0].message.content or ""
        except Exception:
            # very old SDK path
            import openai as old
            resp = old.ChatCompletion.create(model="gpt-4o-mini", messages=messages, temperature=0.0)
            return resp["choices"][0]["message"]["content"]
    except Exception:
        return ""

def extract_pdf_text(pdf_bytes: bytes, max_pages: int = 50) -> str:
    reader = PdfReader(BytesIO(pdf_bytes))
    text_parts: List[str] = []
    for i, page in enumerate(reader.pages[:max_pages]):
        try:
            text_parts.append(page.extract_text() or "")
        except Exception:
            pass
    return "\n\n".join([t for t in text_parts if t.strip()])

def extract_pptx_text(pptx_bytes: bytes) -> str:
    prs = Presentation(BytesIO(pptx_bytes))
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n\n".join(texts)

def extract_image_text(img_bytes: bytes) -> str:
    # OCR via Vision
    return _vision_extract_text(img_bytes)

def extract_any(files) -> str:
    """
    files: list[UploadedFile] from streamlit.file_uploader(accept_multiple_files=True)
    Returns concatenated text from all files.
    """
    parts: List[str] = []
    for f in files:
        name = (f.name or "").lower()
        data = f.read()
        if name.endswith(".pdf"):
            parts.append(extract_pdf_text(data, max_pages=50))
        elif name.endswith(".pptx"):
            parts.append(extract_pptx_text(data))
        elif name.endswith((".png", ".jpg", ".jpeg")):
            parts.append(extract_image_text(data))
        else:
            # fallback: try decode as text
            try:
                parts.append(data.decode("utf-8", errors="ignore"))
            except Exception:
                pass
    return "\n\n".join([p for p in parts if p.strip()])
